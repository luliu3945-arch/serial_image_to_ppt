from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image, ImageChops, ImageDraw, ImageStat
from pptx import Presentation


def main() -> None:
    parser = argparse.ArgumentParser(description="Generate visual QA evidence and validate one refined editable slide.")
    parser.add_argument("--page", type=int, required=True)
    parser.add_argument("--source", type=Path, required=True)
    parser.add_argument("--preview", type=Path, required=True)
    parser.add_argument("--pptx", type=Path, required=True)
    parser.add_argument("--bundle", type=Path, required=True)
    args = parser.parse_args()
    manifest_path = args.bundle / "visual_elements_manifest.json"
    for item in (args.source, args.preview, args.pptx, manifest_path):
        if not item.exists():
            raise SystemExit(f"Missing required input: {item}")
    rows = json.loads(manifest_path.read_text(encoding="utf-8"))
    source = Image.open(args.source).convert("RGBA")
    overlay = source.copy()
    draw = ImageDraw.Draw(overlay)
    for row in rows:
        x, y, w, h = map(float, (row["x"], row["y"], row["width"], row["height"]))
        color = (0, 90, 210, 210) if row.get("is_text_only") == "yes" else (225, 75, 35, 200)
        draw.rectangle((x, y, x + w, y + h), outline=color, width=2)
    review = args.bundle / "review"
    preview_dir = args.bundle / "quality_preview"
    report_dir = args.bundle / "quality_report"
    review.mkdir(parents=True, exist_ok=True)
    preview_dir.mkdir(parents=True, exist_ok=True)
    report_dir.mkdir(parents=True, exist_ok=True)
    overlay.save(review / f"page_{args.page}_refined_elements_overlay.png")
    rendered = Image.open(args.preview).convert("RGB")
    reference = Image.open(args.source).convert("RGB").resize(rendered.size, Image.Resampling.LANCZOS)
    diff = ImageChops.difference(rendered, reference)
    mad = sum(ImageStat.Stat(diff).mean) / 3
    rendered.save(preview_dir / f"page_{args.page}_refined_preview.png")
    reference.save(report_dir / f"page_{args.page}_original.png")
    rendered.save(report_dir / f"page_{args.page}_recomposed.png")
    diff.save(report_dir / f"page_{args.page}_diff.png")
    deck = Presentation(args.pptx)
    slides = len(deck.slides)
    shapes = list(deck.slides[0].shapes) if slides else []
    text_count = sum(1 for shape in shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
    picture_count = sum(1 for shape in shapes if shape.shape_type == 13)
    width, height = source.size
    bounds = [row.get("element_id") for row in rows if float(row.get("x", 0)) < 0 or float(row.get("y", 0)) < 0 or float(row.get("x", 0)) + float(row.get("width", 0)) > width or float(row.get("y", 0)) + float(row.get("height", 0)) > height]
    passed = slides == 1 and len(shapes) == len(rows) and text_count > 0 and picture_count > 0 and not bounds
    report = {
        "slide_id": f"page_{args.page}",
        "status": "rendered_and_reviewed" if passed else "failed",
        "mean_absolute_rgb_diff": round(mad, 3),
        "slide_count": slides,
        "shape_count": len(shapes),
        "manifest_count": len(rows),
        "text_count": text_count,
        "picture_count": picture_count,
        "bounds_issues": bounds,
    }
    (report_dir / "quality_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    print(json.dumps(report, ensure_ascii=False))
    if not passed:
        raise SystemExit(1)


if __name__ == "__main__":
    main()

