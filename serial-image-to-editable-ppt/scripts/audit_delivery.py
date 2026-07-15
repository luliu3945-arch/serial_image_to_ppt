from __future__ import annotations

import argparse
import csv
import json
import zipfile
from pathlib import Path

from pptx import Presentation


def pages(start: int, end: int) -> list[int]:
    if start > end:
        raise SystemExit("Delivery audit requires ascending order: --start must be <= --end")
    return list(range(start, end + 1))


def inspect(output_dir: Path, page: int) -> dict[str, object]:
    pptx = output_dir / f"page_{page}_refined_editable.pptx"
    bundle = output_dir / f"page_{page}_refined_editable_output"
    required = [
        pptx,
        bundle / "visual_elements_manifest.json",
        bundle / "visual_elements_manifest.csv",
        bundle / "baseline_visual_elements_manifest.json",
        bundle / "baseline_visual_elements_manifest.csv",
        bundle / "split_png_elements.zip",
        bundle / "review" / f"page_{page}_refined_elements_overlay.png",
        bundle / "quality_preview" / f"page_{page}_refined_preview.png",
        bundle / "quality_report" / "quality_report.json",
        bundle / "quality_report" / f"page_{page}_original.png",
        bundle / "quality_report" / f"page_{page}_recomposed.png",
        bundle / "quality_report" / f"page_{page}_diff.png",
    ]
    missing = [str(item) for item in required if not item.exists() or item.stat().st_size == 0]
    slide_count = shape_count = text_count = picture_count = manifest_count = bounds_issues = 0
    qa_status = "missing"
    if pptx.exists():
        deck = Presentation(pptx)
        slide_count = len(deck.slides)
        shapes = list(deck.slides[0].shapes) if slide_count else []
        shape_count = len(shapes)
        text_count = sum(1 for shape in shapes if getattr(shape, "has_text_frame", False) and shape.text.strip())
        picture_count = sum(1 for shape in shapes if shape.shape_type == 13)
    manifest_path = bundle / "visual_elements_manifest.json"
    if manifest_path.exists():
        manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
        manifest_count = len(manifest)
        bounds_issues = sum(1 for row in manifest if float(row.get("x", 0)) < 0 or float(row.get("y", 0)) < 0 or float(row.get("x", 0)) + float(row.get("width", 0)) > 1672 or float(row.get("y", 0)) + float(row.get("height", 0)) > 941)
    quality_path = bundle / "quality_report" / "quality_report.json"
    if quality_path.exists():
        qa_status = json.loads(quality_path.read_text(encoding="utf-8")).get("status", "unknown")
    passed = not missing and slide_count == 1 and shape_count > 0 and text_count > 0 and picture_count > 0 and shape_count == manifest_count and bounds_issues == 0 and qa_status == "rendered_and_reviewed"
    return {"page": page, "pptx_bytes": pptx.stat().st_size if pptx.exists() else 0, "slide_count": slide_count, "shape_count": shape_count, "text_count": text_count, "picture_count": picture_count, "manifest_count": manifest_count, "bounds_issues": bounds_issues, "qa_status": qa_status, "missing_artifacts": missing, "passed": passed}


def main() -> None:
    parser = argparse.ArgumentParser(description="Audit and optionally package a serial editable-PPT delivery range.")
    parser.add_argument("--output-dir", type=Path, required=True)
    parser.add_argument("--start", type=int, required=True)
    parser.add_argument("--end", type=int, required=True)
    parser.add_argument("--create-zip", action="store_true")
    args = parser.parse_args()
    page_list = pages(args.start, args.end)
    rows = [inspect(args.output_dir, page) for page in page_list]
    stem = f"pages_{args.start}_to_{args.end}_delivery_audit"
    json_path = args.output_dir / f"{stem}.json"
    csv_path = args.output_dir / f"{stem}.csv"
    summary = {"start": args.start, "end": args.end, "expected_pages": len(page_list), "passed_pages": sum(1 for row in rows if row["passed"]), "all_passed": all(row["passed"] for row in rows), "pages": rows}
    json_path.write_text(json.dumps(summary, ensure_ascii=False, indent=2), encoding="utf-8")
    with csv_path.open("w", newline="", encoding="utf-8-sig") as handle:
        writer = csv.DictWriter(handle, fieldnames=list(rows[0].keys()))
        writer.writeheader()
        for row in rows:
            item = dict(row); item["missing_artifacts"] = ";".join(row["missing_artifacts"]); writer.writerow(item)
    zip_path = None
    if args.create_zip and summary["all_passed"]:
        zip_path = args.output_dir / f"pages_{args.start}_to_{args.end}_refined_editable.zip"
        with zipfile.ZipFile(zip_path, "w", zipfile.ZIP_DEFLATED) as archive:
            for page in page_list:
                pptx = args.output_dir / f"page_{page}_refined_editable.pptx"
                archive.write(pptx, pptx.name)
            archive.write(json_path, json_path.name); archive.write(csv_path, csv_path.name)
    print(json.dumps({"all_passed": summary["all_passed"], "passed_pages": summary["passed_pages"], "expected_pages": len(page_list), "audit_json": str(json_path), "audit_csv": str(csv_path), "zip": str(zip_path) if zip_path else None}, ensure_ascii=False))
    if not summary["all_passed"]:
        raise SystemExit(1)


if __name__ == "__main__":
    main()
