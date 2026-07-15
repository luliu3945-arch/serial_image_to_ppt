#!/usr/bin/env python3
"""Run deterministic parallel batches for codeximage-to-editable-ppt-v1.

This wrapper prevents batch processing from depending on conversation memory.
By default it groups inputs two at a time, writes one fixed batch_instruction.md
per batch, and uses a global worker pool to process item commands in parallel.
The default scheduler activates two batches at once.
"""

from __future__ import annotations

import argparse
import concurrent.futures as futures
import csv
import json
import re
import shutil
import subprocess
import sys
import time
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Dict, Iterable, List, Optional, Sequence, Tuple

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Pt


DEFAULT_PATTERNS = "*.pptx;*.ppt;*.png;*.jpg;*.jpeg;*.webp;*.bmp;*.tif;*.tiff"
DEFAULT_BATCH_SIZE = 2
DEFAULT_BATCH_WORKERS = 2
EMU_PER_INCH = 914400
DEFAULT_SLIDE_W_IN = 13.3333333333


@dataclass
class BatchPlan:
    batch_index: int
    batch_id: str
    input_paths: List[Path]
    batch_dir: Path
    output_dirs: List[Path]
    instruction_path: Path
    log_path: Path
    commands: List[List[str]]
    started_at: str = ""


@dataclass
class ItemResult:
    batch_id: str
    item_index: int
    input_path: str
    output_dir: str
    command: str
    status: str
    return_code: int
    started_at: str
    finished_at: str
    duration_seconds: float
    output: str


@dataclass
class BatchItem:
    batch_index: int
    batch_id: str
    input_count: int
    input_paths: str
    batch_dir: str
    output_dirs: str
    instruction_path: str
    log_path: str
    status: str
    return_code: Optional[int]
    completed_count: int
    failed_count: int
    started_at: str
    finished_at: str
    duration_seconds: float


def now_text() -> str:
    return time.strftime("%Y-%m-%d %H:%M:%S")


def split_patterns(value: str) -> List[str]:
    out: List[str] = []
    for part in re.split(r"[;,]", value):
        part = part.strip()
        if part:
            out.append(part)
    return out or split_patterns(DEFAULT_PATTERNS)


def safe_name(value: str, max_len: int = 80) -> str:
    value = re.sub(r"[^A-Za-z0-9._-]+", "_", value.strip())
    value = value.strip("._-") or "input"
    return value[:max_len]


def iter_dir_files(path: Path, patterns: Sequence[str], recursive: bool) -> Iterable[Path]:
    for pattern in patterns:
        yield from (path.rglob(pattern) if recursive else path.glob(pattern))


def collect_inputs(paths: Sequence[Path], input_dirs: Sequence[Path], patterns: Sequence[str], recursive: bool) -> List[Path]:
    found: List[Path] = []
    for raw in list(paths) + list(input_dirs):
        p = raw.expanduser().resolve()
        if p.is_file():
            found.append(p)
        elif p.is_dir():
            found.extend(sorted(f.resolve() for f in iter_dir_files(p, patterns, recursive) if f.is_file()))
        else:
            raise FileNotFoundError(p)
    dedup: List[Path] = []
    seen = set()
    for p in found:
        key = str(p).lower()
        if key not in seen:
            dedup.append(p)
            seen.add(key)
    return dedup


def resolve_decompose_script(path_arg: Optional[Path]) -> Path:
    if path_arg:
        return path_arg.expanduser().resolve()
    return Path(__file__).resolve().with_name("decompose_visual_elements.py")


def common_decompose_args(args: argparse.Namespace) -> List[str]:
    cmd: List[str] = []
    if args.config:
        cmd += ["--config", str(args.config.expanduser().resolve())]
    cmd += ["--dpi", str(args.dpi)]
    cmd += ["--granularity", args.granularity]
    if args.ocr:
        cmd.append("--ocr")
    cmd += ["--ocr-lang", args.ocr_lang]
    cmd += ["--ocr-confidence-threshold", str(args.ocr_confidence_threshold)]
    if args.editable_text:
        cmd.append("--editable-text")
    cmd += ["--default-font-family", args.default_font_family]
    if args.review:
        cmd.append("--review")
    if args.quality_check:
        cmd.append("--quality-check")
    if args.no_background:
        cmd.append("--no-background")
    if args.clean:
        cmd.append("--clean")
    return cmd


def quote_cmd(parts: Sequence[str]) -> str:
    return subprocess.list2cmdline([str(p) for p in parts])


def chunks(items: Sequence[Path], size: int) -> Iterable[List[Path]]:
    if size < 1:
        raise ValueError("--batch-size must be at least 1")
    for start in range(0, len(items), size):
        yield list(items[start:start + size])


def write_instruction(
    plan: BatchPlan,
    *,
    common_args: Sequence[str],
    batch_size: int,
    batch_workers: int,
    item_workers: int,
    status: str,
    started_at: str,
    finished_at: str = "",
    return_code: Optional[int] = None,
    note: str = "",
) -> None:
    plan.instruction_path.parent.mkdir(parents=True, exist_ok=True)
    input_lines: List[str] = []
    command_lines: List[str] = []
    output_lines: List[str] = []
    for idx, (input_path, output_dir, command) in enumerate(zip(plan.input_paths, plan.output_dirs, plan.commands), start=1):
        input_lines.append(f"{idx}. Source: `{input_path}`")
        input_lines.append(f"   Output: `{output_dir}`")
        output_lines.append(f"{idx}. `{output_dir}`")
        command_lines.append(f"# Item {idx}: {input_path.name}")
        command_lines.append(quote_cmd(command))

    lines = [
        f"# Batch Instruction: {plan.batch_id}",
        "",
        "This file is generated by scripts/run_batches.py. Do not rely on chat memory for this batch.",
        "Batch execution is deterministic: this batch uses the shared parameters below for every item.",
        "",
        "## Fixed Batch Rules",
        "",
        "- Process the inputs listed in this file only.",
        "- Use exactly the shared parameters shown below for every item in this batch.",
        "- Run each item command separately; do not merge outputs from different inputs.",
        "- Keep one output folder per input inside this batch folder.",
        "- Execute through the global worker pool; do not run ad hoc manual loops from chat.",
        "- If one item fails, record the failure in `batch_run.log` and `batch_summary.csv`.",
        "- Do not infer missing batch parameters from chat history.",
        "",
        "## Parallel Execution",
        "",
        f"- Batch size: `{batch_size}`",
        f"- Max active batches: `{batch_workers}`",
        f"- Max concurrent item workers: `{item_workers}`",
        "",
        "## Batch Inputs",
        "",
        *input_lines,
        "",
        "## Shared Parameters",
        "",
        "```text",
        quote_cmd(common_args) if common_args else "(none)",
        "```",
        "",
        "## Commands",
        "",
        "```text",
        "\n".join(command_lines),
        "```",
        "",
        "## Expected Outputs",
        "",
        *output_lines,
        "",
        "Each item output folder should contain:",
        "",
        "- `source_pages/`",
        "- `split_png_elements/`",
        "- `split_png_elements.zip`",
        "- `visual_elements_manifest.csv`",
        "- `visual_elements_manifest.json`",
        "- `recomposed_from_elements.pptx`",
        "- optional `review/` and `quality_report/` when enabled",
        "",
        "After every item in every batch succeeds, the batch root should contain:",
        "",
        "- `merged_recomposed_from_elements.pptx`",
        "- `merge_report.csv`",
        "- `merge_report.json`",
        "",
        "The merged deck order must follow the original input order, not parallel completion order.",
        "",
        "## Status",
        "",
        f"- Status: `{status}`",
        f"- Started at: `{started_at}`",
        f"- Finished at: `{finished_at}`",
        f"- Return code: `{'' if return_code is None else return_code}`",
    ]
    if note:
        lines += ["", "## Note", "", note]
    plan.instruction_path.write_text("\n".join(lines) + "\n", encoding="utf-8")


def build_plans(
    *,
    inputs: Sequence[Path],
    outdir: Path,
    batch_size: int,
    decompose_script: Path,
    python_exe: str,
    common_args: Sequence[str],
) -> List[BatchPlan]:
    plans: List[BatchPlan] = []
    for idx, batch_inputs in enumerate(chunks(inputs, batch_size), start=1):
        batch_id = f"batch_{idx:03d}_{len(batch_inputs)}_items"
        batch_dir = outdir / batch_id
        output_dirs = [
            batch_dir / f"item_{item_idx:02d}_{safe_name(input_path.stem)}"
            for item_idx, input_path in enumerate(batch_inputs, start=1)
        ]
        commands = [
            [python_exe, str(decompose_script), str(input_path), "--outdir", str(output_dir), *common_args]
            for input_path, output_dir in zip(batch_inputs, output_dirs)
        ]
        plans.append(
            BatchPlan(
                batch_index=idx,
                batch_id=batch_id,
                input_paths=batch_inputs,
                batch_dir=batch_dir,
                output_dirs=output_dirs,
                instruction_path=batch_dir / "batch_instruction.md",
                log_path=batch_dir / "batch_run.log",
                commands=commands,
            )
        )
    return plans


def run_item(plan: BatchPlan, item_index: int, dry_run: bool) -> ItemResult:
    input_path = plan.input_paths[item_index - 1]
    output_dir = plan.output_dirs[item_index - 1]
    command = plan.commands[item_index - 1]
    started = now_text()
    t0 = time.time()
    if dry_run:
        return ItemResult(
            batch_id=plan.batch_id,
            item_index=item_index,
            input_path=str(input_path),
            output_dir=str(output_dir),
            command=quote_cmd(command),
            status="dry_run",
            return_code=0,
            started_at=started,
            finished_at=now_text(),
            duration_seconds=round(time.time() - t0, 3),
            output="Dry run: command was not executed.\n",
        )
    proc = subprocess.run(command, text=True, stdout=subprocess.PIPE, stderr=subprocess.STDOUT)
    return ItemResult(
        batch_id=plan.batch_id,
        item_index=item_index,
        input_path=str(input_path),
        output_dir=str(output_dir),
        command=quote_cmd(command),
        status="completed" if proc.returncode == 0 else "failed",
        return_code=proc.returncode,
        started_at=started,
        finished_at=now_text(),
        duration_seconds=round(time.time() - t0, 3),
        output=proc.stdout,
    )


def finalize_batch(
    plan: BatchPlan,
    results: Sequence[ItemResult],
    *,
    common_args: Sequence[str],
    batch_size: int,
    batch_workers: int,
    item_workers: int,
    dry_run: bool,
) -> BatchItem:
    ordered = sorted(results, key=lambda r: r.item_index)
    completed_count = sum(1 for r in ordered if r.return_code == 0)
    failed_count = sum(1 for r in ordered if r.return_code != 0)
    status = "dry_run" if dry_run else ("completed" if failed_count == 0 else "failed")
    return_code = 0 if failed_count == 0 else next(r.return_code for r in ordered if r.return_code != 0)
    finished = now_text()
    started_values = [r.started_at for r in ordered] or [plan.started_at or finished]
    started = plan.started_at or min(started_values)

    log_parts: List[str] = []
    for result in ordered:
        log_parts.append(f"===== Item {result.item_index}: {result.input_path} =====")
        log_parts.append(result.command)
        log_parts.append(f"Started: {result.started_at}")
        log_parts.append(f"Finished: {result.finished_at}")
        log_parts.append(f"Duration seconds: {result.duration_seconds}")
        log_parts.append(f"Return code: {result.return_code}")
        log_parts.append(result.output)
        log_parts.append("")
    plan.log_path.write_text("\n".join(log_parts), encoding="utf-8", errors="replace")

    note = ""
    if failed_count:
        note = f"{failed_count} item(s) failed. See `{plan.log_path}` for decompose script output."
    write_instruction(
        plan,
        common_args=common_args,
        batch_size=batch_size,
        batch_workers=batch_workers,
        item_workers=item_workers,
        status=status,
        started_at=started,
        finished_at=finished,
        return_code=return_code,
        note=note,
    )
    return BatchItem(
        batch_index=plan.batch_index,
        batch_id=plan.batch_id,
        input_count=len(plan.input_paths),
        input_paths=json.dumps([str(p) for p in plan.input_paths], ensure_ascii=False),
        batch_dir=str(plan.batch_dir),
        output_dirs=json.dumps([str(p) for p in plan.output_dirs], ensure_ascii=False),
        instruction_path=str(plan.instruction_path),
        log_path=str(plan.log_path),
        status=status,
        return_code=return_code,
        completed_count=completed_count,
        failed_count=failed_count,
        started_at=started,
        finished_at=finished,
        duration_seconds=round(sum(r.duration_seconds for r in ordered), 3),
    )


def write_summary(items: Sequence[BatchItem], outdir: Path) -> None:
    ordered = sorted(items, key=lambda item: item.batch_index)
    rows = [asdict(item) for item in ordered]
    outdir.mkdir(parents=True, exist_ok=True)
    (outdir / "batch_summary.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")
    if rows:
        with (outdir / "batch_summary.csv").open("w", newline="", encoding="utf-8-sig") as f:
            writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)


def as_int(value: object, default: int = 0) -> int:
    try:
        return int(float(value))  # handles JSON strings and CSV-like values
    except Exception:
        return default


def as_float(value: object, default: float = 0.0) -> float:
    try:
        return float(value)
    except Exception:
        return default


def hex_to_rgb(value: object) -> Optional[Tuple[int, int, int]]:
    if not isinstance(value, str) or not re.match(r"^#[0-9A-Fa-f]{6}$", value):
        return None
    value = value.lstrip("#")
    return int(value[0:2], 16), int(value[2:4], 16), int(value[4:6], 16)


def unique_slide_ids(records: Sequence[dict]) -> List[str]:
    out: List[str] = []
    seen = set()
    for rec in records:
        slide_id = str(rec.get("slide_id") or rec.get("image_id") or "")
        if slide_id and slide_id not in seen:
            out.append(slide_id)
            seen.add(slide_id)
    return out


def set_slide_background(slide, records: Sequence[dict]) -> None:
    for rec in records:
        if rec.get("element_type") != "background":
            continue
        rgb = hex_to_rgb(rec.get("dominant_color"))
        if not rgb:
            return
        slide.background.fill.solid()
        slide.background.fill.fore_color.rgb = RGBColor(*rgb)
        return


def add_merged_textbox(slide, rec: dict, slide_width: int, slide_height: int) -> None:
    source_w = max(1, as_int(rec.get("source_width"), 1))
    source_h = max(1, as_int(rec.get("source_height"), 1))
    left = int(as_int(rec.get("x")) / source_w * slide_width)
    top = int(as_int(rec.get("y")) / source_h * slide_height)
    width = max(1, int(as_int(rec.get("width"), 1) / source_w * slide_width))
    height = max(1, int(as_int(rec.get("height"), 1) / source_h * slide_height))
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    lines = str(rec.get("recognized_text") or "").split("\n") or [""]
    for idx, line in enumerate(lines):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.alignment = PP_ALIGN.CENTER if rec.get("text_alignment_guess") == "center" else PP_ALIGN.LEFT
        run = para.add_run()
        run.text = line
        font = run.font
        font.name = str(rec.get("font_family_guess") or "Arial")
        font.size = Pt(max(4.0, as_float(rec.get("font_size_guess"), 10.0)))
        rgb = hex_to_rgb(rec.get("font_color_guess"))
        font.color.rgb = RGBColor(*(rgb or (0, 0, 0)))
        if rec.get("element_type") == "title_txt":
            font.bold = True


def add_merged_picture(slide, rec: dict, item_output_dir: Path, slide_width: int, slide_height: int) -> bool:
    source_w = max(1, as_int(rec.get("source_width"), 1))
    source_h = max(1, as_int(rec.get("source_height"), 1))
    rel = rec.get("relative_path")
    if not rel:
        return False
    img_path = item_output_dir / str(rel)
    if not img_path.exists():
        return False
    left = int(as_int(rec.get("x")) / source_w * slide_width)
    top = int(as_int(rec.get("y")) / source_h * slide_height)
    width = max(1, int(as_int(rec.get("width"), 1) / source_w * slide_width))
    height = max(1, int(as_int(rec.get("height"), 1) / source_h * slide_height))
    slide.shapes.add_picture(str(img_path), left, top, width=width, height=height)
    return True


def collect_merge_entries(plans: Sequence[BatchPlan]) -> List[dict]:
    entries: List[dict] = []
    original_index = 0
    for plan in sorted(plans, key=lambda p: p.batch_index):
        for item_index, (input_path, output_dir) in enumerate(zip(plan.input_paths, plan.output_dirs), start=1):
            original_index += 1
            pptx_path = output_dir / "recomposed_from_elements.pptx"
            manifest_path = output_dir / "visual_elements_manifest.json"
            entries.append(
                {
                    "original_index": original_index,
                    "batch_index": plan.batch_index,
                    "batch_id": plan.batch_id,
                    "item_index": item_index,
                    "input_path": str(input_path),
                    "output_dir": str(output_dir),
                    "item_pptx": str(pptx_path),
                    "manifest": str(manifest_path),
                    "item_pptx_exists": pptx_path.exists(),
                    "manifest_exists": manifest_path.exists(),
                    "included": False,
                    "slide_count": 0,
                    "note": "",
                }
            )
    return entries


def write_merge_report(entries: Sequence[dict], outdir: Path, merged_path: Optional[Path], status: str, note: str = "") -> None:
    rows = [dict(entry) for entry in entries]
    report = {
        "status": status,
        "merged_pptx": str(merged_path) if merged_path else "",
        "note": note,
        "items": rows,
    }
    (outdir / "merge_report.json").write_text(json.dumps(report, ensure_ascii=False, indent=2), encoding="utf-8")
    if rows:
        with (outdir / "merge_report.csv").open("w", newline="", encoding="utf-8-sig") as f:
            writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            writer.writeheader()
            writer.writerows(rows)


def merge_recomposed_outputs(plans: Sequence[BatchPlan], outdir: Path, output_name: str) -> Path:
    entries = collect_merge_entries(plans)
    missing = [e for e in entries if not e["item_pptx_exists"] or not e["manifest_exists"]]
    if missing:
        note = f"Cannot merge: {len(missing)} item(s) are missing recomposed_from_elements.pptx or visual_elements_manifest.json."
        write_merge_report(entries, outdir, None, "failed", note)
        raise RuntimeError(note)

    item_records: List[Tuple[dict, List[dict]]] = []
    first_dims: Optional[Tuple[int, int]] = None
    for entry in entries:
        records = json.loads(Path(entry["manifest"]).read_text(encoding="utf-8"))
        if not isinstance(records, list):
            raise RuntimeError(f"Manifest is not a list: {entry['manifest']}")
        slide_ids = unique_slide_ids(records)
        entry["slide_count"] = len(slide_ids)
        entry["included"] = True
        for rec in records:
            source_w = as_int(rec.get("source_width"))
            source_h = as_int(rec.get("source_height"))
            if source_w > 0 and source_h > 0:
                first_dims = first_dims or (source_w, source_h)
                break
        item_records.append((entry, records))

    if not first_dims:
        note = "Cannot merge: no source dimensions found in item manifests."
        write_merge_report(entries, outdir, None, "failed", note)
        raise RuntimeError(note)

    source_w, source_h = first_dims
    slide_w_in = DEFAULT_SLIDE_W_IN
    slide_h_in = slide_w_in * source_h / max(1, source_w)
    prs = Presentation()
    prs.slide_width = int(slide_w_in * EMU_PER_INCH)
    prs.slide_height = int(slide_h_in * EMU_PER_INCH)
    blank = prs.slide_layouts[6]

    for entry, records in item_records:
        item_output_dir = Path(entry["output_dir"])
        for slide_id in unique_slide_ids(records):
            slide_records = [r for r in records if str(r.get("slide_id") or r.get("image_id") or "") == slide_id]
            slide = prs.slides.add_slide(blank)
            set_slide_background(slide, slide_records)
            for rec in sorted(slide_records, key=lambda r: as_int(r.get("z_order"))):
                if rec.get("element_type") == "background":
                    continue
                if rec.get("is_text_only") == "yes" and rec.get("converted_to_editable_text") == "yes":
                    add_merged_textbox(slide, rec, prs.slide_width, prs.slide_height)
                else:
                    add_merged_picture(slide, rec, item_output_dir, prs.slide_width, prs.slide_height)

    merged_path = outdir / output_name
    prs.save(merged_path)
    write_merge_report(entries, outdir, merged_path, "completed", f"Merged {sum(e['slide_count'] for e in entries)} slide(s).")
    return merged_path


def run_plans(
    *,
    plans: Sequence[BatchPlan],
    outdir: Path,
    common_args: Sequence[str],
    batch_size: int,
    batch_workers: int,
    item_workers: int,
    dry_run: bool,
    continue_on_error: bool,
) -> List[BatchItem]:
    for plan in plans:
        plan.batch_dir.mkdir(parents=True, exist_ok=True)
        write_instruction(
            plan,
            common_args=common_args,
            batch_size=batch_size,
            batch_workers=batch_workers,
            item_workers=item_workers,
            status="queued",
            started_at="",
        )

    completed: List[BatchItem] = []
    active: Dict[str, Dict[str, object]] = {}
    pending: Dict[futures.Future[ItemResult], Tuple[str, int]] = {}
    next_plan_index = 0
    stop_launching = False

    def launch_plan(executor: futures.Executor, plan: BatchPlan) -> None:
        plan.started_at = now_text()
        write_instruction(
            plan,
            common_args=common_args,
            batch_size=batch_size,
            batch_workers=batch_workers,
            item_workers=item_workers,
            status="running" if not dry_run else "dry_run",
            started_at=plan.started_at,
        )
        active[plan.batch_id] = {"plan": plan, "remaining": len(plan.input_paths), "results": []}
        for item_index in range(1, len(plan.input_paths) + 1):
            fut = executor.submit(run_item, plan, item_index, dry_run)
            pending[fut] = (plan.batch_id, item_index)

    with futures.ThreadPoolExecutor(max_workers=item_workers) as executor:
        while next_plan_index < len(plans) and len(active) < batch_workers:
            launch_plan(executor, plans[next_plan_index])
            next_plan_index += 1

        while pending:
            done, _ = futures.wait(pending.keys(), return_when=futures.FIRST_COMPLETED)
            for fut in done:
                batch_id, _ = pending.pop(fut)
                result = fut.result()
                state = active[batch_id]
                state["results"].append(result)  # type: ignore[index, union-attr]
                state["remaining"] = int(state["remaining"]) - 1

                if int(state["remaining"]) == 0:
                    plan = state["plan"]
                    assert isinstance(plan, BatchPlan)
                    results = state["results"]
                    assert isinstance(results, list)
                    item = finalize_batch(
                        plan,
                        results,
                        common_args=common_args,
                        batch_size=batch_size,
                        batch_workers=batch_workers,
                        item_workers=item_workers,
                        dry_run=dry_run,
                    )
                    completed.append(item)
                    write_summary(completed, outdir)
                    del active[batch_id]
                    if item.status == "failed" and not continue_on_error:
                        stop_launching = True

            while not stop_launching and next_plan_index < len(plans) and len(active) < batch_workers:
                launch_plan(executor, plans[next_plan_index])
                next_plan_index += 1

    return completed


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(
        description="Run image-PPT decomposition in deterministic parallel batches using one shared parameter set."
    )
    p.add_argument("inputs", nargs="*", type=Path, help="Input files or directories.")
    p.add_argument("--input-dir", action="append", type=Path, default=[], help="Directory containing input files. Repeatable.")
    p.add_argument("--patterns", default=DEFAULT_PATTERNS, help="File globs for directories, separated by semicolon or comma.")
    p.add_argument("--recursive", action="store_true", help="Search input directories recursively.")
    p.add_argument("--outdir", type=Path, default=Path("batch_output"), help="Root output directory for all batches.")
    p.add_argument("--batch-size", type=int, default=DEFAULT_BATCH_SIZE, help="Number of input files per batch. Default: 2.")
    p.add_argument("--batch-workers", type=int, default=DEFAULT_BATCH_WORKERS, help="Maximum active batches at once. Default: 2.")
    p.add_argument("--workers", type=int, help="Maximum concurrent item workers. Default: batch-size * batch-workers.")
    p.add_argument("--python", default=sys.executable, help="Python executable used to run decompose_visual_elements.py.")
    p.add_argument("--decompose-script", type=Path, help="Override path to decompose_visual_elements.py.")
    p.add_argument("--dry-run", action="store_true", help="Write batch instructions without running decomposition.")
    p.add_argument("--continue-on-error", action="store_true", help="Continue launching later batches after a batch failure.")
    p.add_argument("--clean-batch-root", action="store_true", help="Delete the batch root output directory before running.")
    p.add_argument("--no-merge", action="store_true", help="Skip final merged_recomposed_from_elements.pptx generation.")
    p.add_argument("--merge-output", default="merged_recomposed_from_elements.pptx", help="Merged PPTX filename at the batch root.")

    p.add_argument("--config", type=Path, help="Optional YAML/JSON config file passed to every batch.")
    p.add_argument("--dpi", type=int, default=300, help="DPI for fallback rendering.")
    p.add_argument("--granularity", choices=["coarse", "normal", "fine", "ultra"], default="fine")
    p.add_argument("--ocr", action="store_true", help="Enable OCR for every batch.")
    p.add_argument("--ocr-lang", default="eng", help="Tesseract language, e.g. chi_sim+eng.")
    p.add_argument("--ocr-confidence-threshold", type=float, default=70.0)
    p.add_argument("--editable-text", action="store_true", help="Convert reliable text-only PNG elements into editable text boxes.")
    p.add_argument("--default-font-family", default="Arial")
    p.add_argument("--review", action="store_true", help="Generate review overlays for every batch.")
    p.add_argument("--quality-check", action="store_true", help="Generate diff reports for every batch.")
    p.add_argument("--no-background", action="store_true", help="Do not export a solid background PNG element.")
    p.add_argument("--clean", action="store_true", help="Pass --clean to each decompose run.")
    return p.parse_args()


def main() -> int:
    args = parse_args()
    if args.batch_size < 1:
        raise ValueError("--batch-size must be at least 1")
    if args.batch_workers < 1:
        raise ValueError("--batch-workers must be at least 1")
    item_workers = args.workers if args.workers is not None else args.batch_size * args.batch_workers
    if item_workers < 1:
        raise ValueError("--workers must be at least 1")

    patterns = split_patterns(args.patterns)
    inputs = collect_inputs(args.inputs, args.input_dir, patterns, args.recursive)
    if not inputs:
        print("No input files found.", file=sys.stderr)
        return 2

    decompose_script = resolve_decompose_script(args.decompose_script)
    if not decompose_script.exists():
        raise FileNotFoundError(decompose_script)

    outdir = args.outdir.expanduser().resolve()
    if args.clean_batch_root and outdir.exists():
        shutil.rmtree(outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    common_args = common_decompose_args(args)
    plans = build_plans(
        inputs=inputs,
        outdir=outdir,
        batch_size=args.batch_size,
        decompose_script=decompose_script,
        python_exe=args.python,
        common_args=common_args,
    )
    print(
        f"Prepared {len(plans)} batch(es), batch-size={args.batch_size}, "
        f"batch-workers={args.batch_workers}, item-workers={item_workers}."
    )

    completed = run_plans(
        plans=plans,
        outdir=outdir,
        common_args=common_args,
        batch_size=args.batch_size,
        batch_workers=args.batch_workers,
        item_workers=item_workers,
        dry_run=args.dry_run,
        continue_on_error=args.continue_on_error,
    )
    write_summary(completed, outdir)

    failed = [item for item in completed if item.status == "failed"]
    not_started = len(plans) - len(completed)
    if failed:
        print(
            f"Completed with {len(failed)} failed batch(es); {not_started} batch(es) not launched. "
            f"Summary: {outdir / 'batch_summary.csv'}",
            file=sys.stderr,
        )
        return 1
    if not args.dry_run and not args.no_merge:
        merged = merge_recomposed_outputs(plans, outdir, args.merge_output)
        print(f"Merged PPTX: {merged}")
    print(f"Completed {len(completed)} batch(es). Summary: {outdir / 'batch_summary.csv'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
