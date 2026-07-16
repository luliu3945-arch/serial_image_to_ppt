#!/usr/bin/env python3
"""codeximage-to-editable-ppt-v1

High-fidelity baseline utility for image-based PPT/PPTX and slide screenshots.
It extracts or renders raster slide pages, decomposes them into independent PNG
visual elements, writes CSV/JSON manifests, and recomposes a PPTX using PNG
assets plus editable text boxes for OCR-confident text-only elements.

This tool is intentionally image-centric. It does not attempt to recover native
PowerPoint vector objects from a rasterized slide.
"""

from __future__ import annotations

import argparse
import csv
import json
import math
import os
import re
import shutil
import subprocess
import sys
import tempfile
import zipfile
from dataclasses import asdict, dataclass
from pathlib import Path
from typing import Any, Dict, Iterable, List, Optional, Tuple

import cv2
import numpy as np
from PIL import Image, ImageChops, ImageDraw, ImageFont
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    import yaml
except Exception:  # pragma: no cover
    yaml = None

EMU_PER_INCH = 914400
DEFAULT_SLIDE_W_IN = 13.3333333333
DEFAULT_SLIDE_H_IN = 7.5
TRANSPARENT_REQUIRED_TYPES = {
    "text",
    "icon",
    "logo",
    "symbol",
    "button",
    "badge",
    "arrow",
    "line",
    "divider",
    "shape",
    "simple_shape",
}


@dataclass
class PageSource:
    slide_id: str
    source_path: str
    page_source_type: str
    source_width: int
    source_height: int
    note: str = ""


@dataclass
class ElementRecord:
    slide_id: str
    element_id: str
    file_name: str
    relative_path: str
    element_type: str
    is_text_only: str
    recognized_text: str
    OCR_confidence: float
    x: int
    y: int
    width: int
    height: int
    z_order: int
    opacity: float
    rotation: float
    crop_mask_type: str
    transparent_background: str
    alpha_note: str
    shadow_effect: str
    converted_to_editable_text: str
    conversion_note: str
    page_source_type: str
    source_page_path: str
    source_width: int
    source_height: int
    bbox_pixels: str
    bbox_normalized: str
    pptx_left_emu: int
    pptx_top_emu: int
    pptx_width_emu: int
    pptx_height_emu: int
    dominant_color: str
    font_family_guess: str
    font_size_guess: float
    font_color_guess: str
    text_alignment_guess: str
    language_guess: str
    split_confidence: float
    parent_group_id: str
    needs_manual_review: str


def run(cmd: List[str], cwd: Optional[Path] = None, check: bool = True) -> subprocess.CompletedProcess:
    proc = subprocess.run(
        cmd,
        cwd=str(cwd) if cwd else None,
        text=True,
        stdout=subprocess.PIPE,
        stderr=subprocess.STDOUT,
    )
    if check and proc.returncode != 0:
        raise RuntimeError(f"Command failed: {' '.join(cmd)}\n{proc.stdout}")
    return proc


def libreoffice_command() -> str:
    executable = shutil.which("libreoffice") or shutil.which("soffice")
    if not executable:
        raise RuntimeError("LibreOffice was not found on PATH (expected libreoffice or soffice).")
    return executable


def safe_mkdir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def load_config(path: Optional[Path]) -> Dict[str, Any]:
    if not path:
        return {}
    text = path.read_text(encoding="utf-8")
    if path.suffix.lower() in {".yaml", ".yml"} and yaml:
        return yaml.safe_load(text) or {}
    return json.loads(text)


def apply_config_defaults(args: argparse.Namespace, cfg: Dict[str, Any]) -> argparse.Namespace:
    render = cfg.get("render", {}) or {}
    seg = cfg.get("segmentation", {}) or {}
    ocr = cfg.get("ocr", {}) or {}
    rec = cfg.get("recompose", {}) or {}
    qc = cfg.get("quality_check", {}) or {}

    if "dpi" in render and args.dpi == parser.get_default("dpi"):
        args.dpi = int(render["dpi"])
    if "granularity" in seg and args.granularity == parser.get_default("granularity"):
        args.granularity = str(seg["granularity"])
    if seg.get("export_plain_background") is False:
        args.no_background = True
    if ocr.get("enabled") is True:
        args.ocr = True
    if "language" in ocr and args.ocr_lang == parser.get_default("ocr_lang"):
        args.ocr_lang = str(ocr["language"])
    if "confidence_threshold" in ocr and args.ocr_confidence_threshold == parser.get_default("ocr_confidence_threshold"):
        args.ocr_confidence_threshold = float(ocr["confidence_threshold"])
    if rec.get("editable_text") is True:
        args.editable_text = True
    if "default_font_family" in rec and args.default_font_family == parser.get_default("default_font_family"):
        args.default_font_family = str(rec["default_font_family"])
    if qc.get("generate_review_overlays") is True or qc.get("enabled") is True:
        args.review = True
    if qc.get("generate_diff_images") is True or qc.get("enabled") is True:
        args.quality_check = True
    return args


def page_id_for(index: int, prefix: str = "slide") -> str:
    return f"{prefix}{index:02d}"


def pil_open_rgb(path: Path) -> Image.Image:
    return Image.open(path).convert("RGB")


def save_png_no_compress(img: Image.Image, path: Path) -> None:
    safe_mkdir(path.parent)
    img.save(path, format="PNG", compress_level=0)


def estimate_background_color(img: Image.Image) -> Tuple[int, int, int]:
    """Estimate the dominant plain background color.

    Image-based slides often have white paper plus colored decorative corners.
    A simple corner median can be biased by the decoration. Prefer bright,
    low-saturation border pixels, which usually represent the slide canvas.
    """
    arr = np.array(img.convert("RGB"))
    h, w = arr.shape[:2]
    band = max(6, min(h, w) // 45)
    border = np.concatenate([
        arr[:band, :, :].reshape(-1, 3),
        arr[-band:, :, :].reshape(-1, 3),
        arr[:, :band, :].reshape(-1, 3),
        arr[:, -band:, :].reshape(-1, 3),
    ], axis=0)
    hsv = cv2.cvtColor(border.reshape(-1, 1, 3).astype(np.uint8), cv2.COLOR_RGB2HSV).reshape(-1, 3)
    brightness = border.mean(axis=1)
    saturation = hsv[:, 1]
    candidates = border[(brightness >= np.percentile(brightness, 70)) & (saturation <= 45)]
    if len(candidates) < 50:
        candidates = border[brightness >= np.percentile(brightness, 85)]
    med = np.median(candidates, axis=0).astype(int)
    return tuple(int(v) for v in med.tolist())


def is_picture_shape(shape: Any) -> bool:
    try:
        return shape.shape_type == MSO_SHAPE_TYPE.PICTURE
    except Exception:
        return False


def shape_coverage(shape: Any, slide_w: int, slide_h: int) -> float:
    try:
        x1 = max(0, int(shape.left)); y1 = max(0, int(shape.top))
        x2 = min(slide_w, int(shape.left + shape.width)); y2 = min(slide_h, int(shape.top + shape.height))
        return max(0, x2 - x1) * max(0, y2 - y1) / float(slide_w * slide_h)
    except Exception:
        return 0.0


def extract_full_slide_images_from_pptx(input_path: Path, outdir: Path) -> List[PageSource]:
    """Extract original embedded image when a slide is essentially one full-slide picture.

    If detection fails for a slide, source_path is blank and the caller should render.
    """
    sources: List[PageSource] = []
    prs = Presentation(str(input_path))
    slide_w, slide_h = int(prs.slide_width), int(prs.slide_height)
    safe_mkdir(outdir)

    for idx, slide in enumerate(prs.slides, start=1):
        sid = page_id_for(idx, "slide")
        pictures = [sh for sh in slide.shapes if is_picture_shape(sh)]
        best = None
        best_cov = 0.0
        for pic in pictures:
            cov = shape_coverage(pic, slide_w, slide_h)
            if cov > best_cov:
                best = pic
                best_cov = cov
        if best is not None and best_cov >= 0.88 and len(slide.shapes) <= max(3, len(pictures) + 1):
            blob = best.image.blob
            ext = best.image.ext or "png"
            raw_path = outdir / f"{sid}_embedded.{ext}"
            raw_path.write_bytes(blob)
            try:
                img = Image.open(raw_path).convert("RGBA")
                png_path = outdir / f"{sid}_source.png"
                save_png_no_compress(img, png_path)
                sources.append(PageSource(sid, str(png_path), "full_slide_image_extracted", img.width, img.height, f"coverage={best_cov:.3f}"))
            except Exception as e:
                sources.append(PageSource(sid, "", "render_required", 0, 0, f"embedded extraction failed: {e}"))
        else:
            note = f"pictures={len(pictures)}, shapes={len(slide.shapes)}, best_coverage={best_cov:.3f}"
            sources.append(PageSource(sid, "", "render_required", 0, 0, note))
    return sources


def render_pptx_to_pngs(input_path: Path, outdir: Path, dpi: int) -> List[Path]:
    safe_mkdir(outdir)
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        run([libreoffice_command(), "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(input_path)])
        pdfs = sorted(tmp.glob("*.pdf"))
        if not pdfs:
            raise RuntimeError("LibreOffice did not produce a PDF. Is LibreOffice installed?")
        prefix = outdir / "slide"
        run(["pdftoppm", "-r", str(dpi), "-png", str(pdfs[0]), str(prefix)])
    pngs = sorted(outdir.glob("slide-*.png"))
    if not pngs:
        raise RuntimeError("pdftoppm did not produce page PNG files. Is poppler-utils installed?")
    return pngs


def collect_page_sources(input_path: Path, outdir: Path, dpi: int) -> List[PageSource]:
    ext = input_path.suffix.lower()
    source_dir = outdir / "source_pages"
    safe_mkdir(source_dir)
    sources: List[PageSource] = []

    if ext == ".pptx":
        extracted = extract_full_slide_images_from_pptx(input_path, source_dir)
        if all(src.source_path for src in extracted):
            return extracted
        rendered = render_pptx_to_pngs(input_path, outdir / "rendered_pages", dpi)
        render_by_index = {i + 1: p for i, p in enumerate(rendered)}
        for i, src in enumerate(extracted, start=1):
            if src.source_path:
                sources.append(src)
            else:
                png = render_by_index[i]
                img = pil_open_rgb(png)
                dst = source_dir / f"slide{i:02d}_source.png"
                save_png_no_compress(img, dst)
                sources.append(PageSource(page_id_for(i, "slide"), str(dst), "rendered_slide_image", img.width, img.height, src.note))
        return sources

    if ext == ".ppt":
        rendered = render_pptx_to_pngs(input_path, outdir / "rendered_pages", dpi)
        for i, png in enumerate(rendered, start=1):
            img = pil_open_rgb(png)
            dst = source_dir / f"slide{i:02d}_source.png"
            save_png_no_compress(img, dst)
            sources.append(PageSource(page_id_for(i, "slide"), str(dst), "rendered_slide_image", img.width, img.height, "ppt_rendered"))
        return sources

    if ext in {".png", ".jpg", ".jpeg", ".webp", ".bmp", ".tif", ".tiff"}:
        img = Image.open(input_path).convert("RGBA")
        dst = source_dir / "image01_source.png"
        save_png_no_compress(img, dst)
        return [PageSource("image01", str(dst), "input_image", img.width, img.height, "")]

    raise ValueError(f"Unsupported input file type: {input_path.suffix}")


def write_page_source_report(sources: List[PageSource], outdir: Path) -> None:
    rows = [asdict(s) for s in sources]
    if not rows:
        return
    csv_path = outdir / "image_source_report.csv"
    with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader(); writer.writerows(rows)
    (outdir / "image_source_report.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")


def granularity_params(granularity: str, page_area: int) -> Dict[str, Any]:
    granularity = granularity.lower()
    base = {
        "coarse": dict(min_area=max(120, int(page_area * 0.00045)), close_kernel=(21, 15), dilate=2),
        "normal": dict(min_area=max(60, int(page_area * 0.00018)), close_kernel=(11, 7), dilate=1),
        "fine": dict(min_area=max(24, int(page_area * 0.00007)), close_kernel=(7, 5), dilate=1),
        "ultra": dict(min_area=max(8, int(page_area * 0.000025)), close_kernel=(3, 3), dilate=0),
    }
    return base.get(granularity, base["fine"])


def text_boxes_from_ocr(img: Image.Image, lang: str, enabled: bool, min_conf: float) -> List[Dict[str, Any]]:
    if not enabled or pytesseract is None:
        return []
    rgb = img.convert("RGB")
    config = "--psm 6"
    try:
        data = pytesseract.image_to_data(rgb, lang=lang, output_type=pytesseract.Output.DICT, config=config)
    except Exception:
        try:
            data = pytesseract.image_to_data(rgb, output_type=pytesseract.Output.DICT, config=config)
        except Exception:
            return []

    words: List[Dict[str, Any]] = []
    n = len(data.get("text", []))
    for i in range(n):
        txt = str(data["text"][i]).strip()
        if not txt:
            continue
        try:
            conf = float(data.get("conf", [0])[i])
        except Exception:
            conf = 0.0
        if conf < max(1, min_conf * 0.40):
            continue
        x = int(data["left"][i]); y = int(data["top"][i]); w = int(data["width"][i]); h = int(data["height"][i])
        if w <= 1 or h <= 1:
            continue
        words.append(dict(text=txt, conf=conf, x=x, y=y, w=w, h=h, cx=x+w/2, cy=y+h/2))
    if not words:
        return []

    # Group into lines.
    words.sort(key=lambda d: (d["y"], d["x"]))
    lines: List[List[Dict[str, Any]]] = []
    for word in words:
        placed = False
        for line in lines:
            heights = [w["h"] for w in line]
            median_h = float(np.median(heights))
            line_cy = float(np.median([w["cy"] for w in line]))
            if abs(word["cy"] - line_cy) <= max(6, median_h * 0.55):
                line.append(word); placed = True; break
        if not placed:
            lines.append([word])
    line_objs: List[Dict[str, Any]] = []
    for line in lines:
        line.sort(key=lambda d: d["x"])
        x1 = min(w["x"] for w in line); y1 = min(w["y"] for w in line)
        x2 = max(w["x"] + w["w"] for w in line); y2 = max(w["y"] + w["h"] for w in line)
        text = " ".join(w["text"] for w in line)
        conf = float(np.mean([w["conf"] for w in line]))
        line_objs.append(dict(text=text, conf=conf, x=x1, y=y1, w=x2-x1, h=y2-y1))
    line_objs.sort(key=lambda d: (d["y"], d["x"]))

    # Group nearby lines into text blocks.
    blocks: List[List[Dict[str, Any]]] = []
    for ln in line_objs:
        placed = False
        for block in blocks:
            x1 = min(l["x"] for l in block); x2 = max(l["x"] + l["w"] for l in block)
            y2 = max(l["y"] + l["h"] for l in block)
            med_h = float(np.median([l["h"] for l in block]))
            vertical_gap = ln["y"] - y2
            x_overlap = min(x2, ln["x"] + ln["w"]) - max(x1, ln["x"])
            close_left = abs(ln["x"] - x1) < max(40, 0.20 * max(x2-x1, ln["w"]))
            overlaps = x_overlap > 0.25 * min(max(1, x2-x1), max(1, ln["w"]))
            if 0 <= vertical_gap <= max(12, med_h * 1.35) and (close_left or overlaps):
                block.append(ln); placed = True; break
        if not placed:
            blocks.append([ln])

    out: List[Dict[str, Any]] = []
    for block in blocks:
        block.sort(key=lambda d: (d["y"], d["x"]))
        x1 = min(l["x"] for l in block); y1 = min(l["y"] for l in block)
        x2 = max(l["x"] + l["w"] for l in block); y2 = max(l["y"] + l["h"] for l in block)
        text = "\n".join(l["text"] for l in block)
        conf = float(np.mean([l["conf"] for l in block]))
        out.append(dict(text=text, conf=conf, x=x1, y=y1, w=x2-x1, h=y2-y1, lines=len(block)))
    return out


def non_background_mask(img_rgb: np.ndarray, background_color: Tuple[int, int, int]) -> np.ndarray:
    bg = np.array(background_color, dtype=np.int16)
    arr = img_rgb.astype(np.int16)
    diff_bg = np.max(np.abs(arr - bg), axis=2)
    hsv = cv2.cvtColor(img_rgb, cv2.COLOR_RGB2HSV)
    saturation = hsv[:, :, 1]
    # Detect content on light or colored pages.
    mask = ((diff_bg > 18) | (saturation > 35)).astype(np.uint8) * 255
    mask = cv2.medianBlur(mask, 3)
    return mask


def subtract_text_regions(mask: np.ndarray, text_boxes: List[Dict[str, Any]], pad: int = 3) -> np.ndarray:
    out = mask.copy()
    h, w = out.shape
    for b in text_boxes:
        x1 = max(0, int(b["x"]) - pad); y1 = max(0, int(b["y"]) - pad)
        x2 = min(w, int(b["x"] + b["w"]) + pad); y2 = min(h, int(b["y"] + b["h"]) + pad)
        out[y1:y2, x1:x2] = 0
    return out


def connected_boxes(mask: np.ndarray, min_area: int) -> List[Tuple[int, int, int, int, int]]:
    n, labels, stats, _ = cv2.connectedComponentsWithStats(mask, 8)
    boxes: List[Tuple[int, int, int, int, int]] = []
    h, w = mask.shape
    for i in range(1, n):
        x, y, bw, bh, area = stats[i]
        if int(area) < min_area:
            continue
        if bw < 2 or bh < 2:
            continue
        # Faint gradients or anti-aliased background decorations can connect
        # many objects into one near-full-page component. Ignore that component
        # and let smaller components/background entries carry the layout.
        if (bw >= w * 0.80 and bh >= h * 0.80) or (area > w * h * 0.45):
            continue
        boxes.append((int(x), int(y), int(bw), int(bh), int(area)))
    return boxes


def boxes_iou(a: Tuple[int, int, int, int], b: Tuple[int, int, int, int]) -> float:
    ax, ay, aw, ah = a; bx, by, bw, bh = b
    ax2, ay2 = ax + aw, ay + ah; bx2, by2 = bx + bw, by + bh
    ix1, iy1 = max(ax, bx), max(ay, by)
    ix2, iy2 = min(ax2, bx2), min(ay2, by2)
    inter = max(0, ix2 - ix1) * max(0, iy2 - iy1)
    union = aw * ah + bw * bh - inter
    return inter / union if union else 0.0


def merge_boxes(boxes: List[Tuple[int, int, int, int, int]], mode: str, page_w: int, page_h: int) -> List[Tuple[int, int, int, int, int]]:
    """Pad boxes and optionally merge for coarse mode.

    For image-based slides, over-merging is worse than mild fragmentation.
    Text grouping is handled by OCR, so normal/fine/ultra modes keep visual
    components separate. Only coarse mode merges nearby boxes into large areas.
    """
    pad = {"ultra": 1, "fine": 2, "normal": 3, "coarse": 5}.get(mode, 2)
    if mode != "coarse":
        out = []
        for x, y, w, h, area in boxes:
            x1 = max(0, x - pad); y1 = max(0, y - pad)
            x2 = min(page_w, x + w + pad); y2 = min(page_h, y + h + pad)
            out.append((x1, y1, x2 - x1, y2 - y1, area))
        return sorted(out, key=lambda b: (b[1], b[0], b[2] * b[3]))

    rects = [[x, y, x+w, y+h, a] for x, y, w, h, a in boxes]
    changed = True
    while changed:
        changed = False
        used = [False] * len(rects)
        new: List[List[int]] = []
        for i, r in enumerate(rects):
            if used[i]:
                continue
            x1, y1, x2, y2, area = r
            used[i] = True
            for j in range(i+1, len(rects)):
                if used[j]:
                    continue
                ux1, uy1, ux2, uy2, uarea = rects[j]
                gap_x = max(0, max(ux1 - x2, x1 - ux2))
                gap_y = max(0, max(uy1 - y2, y1 - uy2))
                h1 = y2 - y1; h2 = uy2 - uy1
                if gap_x <= max(20, 0.45 * max(h1, h2)) and gap_y <= max(18, 0.45 * max(h1, h2)):
                    x1, y1 = min(x1, ux1), min(y1, uy1)
                    x2, y2 = max(x2, ux2), max(y2, uy2)
                    area += uarea
                    used[j] = True
                    changed = True
            new.append([x1, y1, x2, y2, area])
        rects = new
    out = []
    for x1, y1, x2, y2, area in rects:
        x1 = max(0, x1 - pad); y1 = max(0, y1 - pad)
        x2 = min(page_w, x2 + pad); y2 = min(page_h, y2 + pad)
        out.append((x1, y1, x2 - x1, y2 - y1, area))
    return sorted(out, key=lambda b: (b[1], b[0], b[2] * b[3]))


def make_alpha_crop(img_rgba: Image.Image, content_mask: np.ndarray, bbox: Tuple[int, int, int, int], soft_edges: bool = True) -> Image.Image:
    x, y, w, h = bbox
    crop = img_rgba.crop((x, y, x+w, y+h)).convert("RGBA")
    mask_crop = content_mask[y:y+h, x:x+w]
    if mask_crop.size == 0:
        return crop
    alpha = mask_crop.copy()
    alpha = cv2.dilate(alpha, cv2.getStructuringElement(cv2.MORPH_RECT, (3, 3)), iterations=1)
    if soft_edges:
        alpha = cv2.GaussianBlur(alpha, (3, 3), 0)
    arr = np.array(crop)
    arr[:, :, 3] = alpha
    return Image.fromarray(arr)


def crop_text_png(img_rgba: Image.Image, bbox: Tuple[int, int, int, int], background_color: Tuple[int, int, int]) -> Image.Image:
    x, y, w, h = bbox
    pad = 3
    x1, y1 = max(0, x-pad), max(0, y-pad)
    x2, y2 = min(img_rgba.width, x+w+pad), min(img_rgba.height, y+h+pad)
    crop = img_rgba.crop((x1, y1, x2, y2)).convert("RGBA")
    rgb = np.array(crop.convert("RGB"))
    bg = np.array(background_color, dtype=np.int16)
    diff = np.max(np.abs(rgb.astype(np.int16) - bg), axis=2)
    hsv = cv2.cvtColor(rgb, cv2.COLOR_RGB2HSV)
    mask = ((diff > 14) | (hsv[:, :, 1] > 25)).astype(np.uint8) * 255
    mask = cv2.dilate(mask, cv2.getStructuringElement(cv2.MORPH_RECT, (2, 2)), iterations=1)
    mask = cv2.GaussianBlur(mask, (3, 3), 0)
    arr = np.array(crop)
    arr[:, :, 3] = mask
    return Image.fromarray(arr)


def dominant_color_hex(img: Image.Image) -> str:
    arr = np.array(img.convert("RGBA"))
    alpha = arr[:, :, 3]
    pix = arr[alpha > 32, :3]
    if pix.size == 0:
        return ""
    med = np.median(pix, axis=0).astype(int)
    return "#%02X%02X%02X" % tuple(med.tolist())


def infer_crop_mask_type(alpha_img: Image.Image) -> str:
    arr = np.array(alpha_img.convert("RGBA"))
    alpha = arr[:, :, 3]
    if alpha.size == 0:
        return "rectangle"
    filled = np.mean(alpha > 20)
    h, w = alpha.shape
    aspect = w / max(h, 1)
    if filled > 0.94:
        return "rectangle"
    if 0.82 <= aspect <= 1.22 and 0.60 <= filled <= 0.88:
        return "circle"
    if 0.45 <= filled <= 0.88:
        return "ellipse" if 0.45 <= aspect <= 2.2 else "irregular_mask"
    return "irregular_mask"


def has_transparent_background(img: Image.Image) -> bool:
    alpha = np.array(img.convert("RGBA"))[:, :, 3]
    if alpha.size == 0:
        return False
    return bool(np.mean(alpha < 250) > 0.001 or alpha.min() < 220)


def requires_transparent_background(element_type: str, is_text_only: str) -> bool:
    if is_text_only == "yes":
        return True
    normalized = element_type.lower().replace("-", "_")
    if normalized.endswith("_txt"):
        return True
    return normalized in TRANSPARENT_REQUIRED_TYPES


def alpha_fields_for(crop: Image.Image, element_type: str, is_text_only: str) -> Tuple[str, str]:
    transparent = "yes" if has_transparent_background(crop) else "no"
    if element_type == "background":
        return "no", "full_background"
    if transparent == "yes":
        return transparent, "alpha_mask_applied"
    if requires_transparent_background(element_type, is_text_only):
        return transparent, "transparent_required_but_unreliable; rectangular_crop_fallback"
    return transparent, "rectangular_crop_fallback"


def infer_text_type(box: Dict[str, Any], page_w: int, page_h: int, text: str) -> str:
    y = box["y"]; h = box["h"]; w = box["w"]
    low = text.lower()
    if y < page_h * 0.22 and (h > page_h * 0.025 or w > page_w * 0.25):
        return "title_txt"
    if "doi" in low:
        return "doi_txt"
    if re.search(r"fig\.?|table|source|note", low):
        return "caption_txt"
    if len(text) <= 30 and (re.search(r"^[\d\.\-–—]+$", text) or re.search(r"[=±×÷≤≥μΔ√∑]", text)):
        return "formula_label_txt"
    if w < page_w * 0.18 and h < page_h * 0.06:
        return "label_txt"
    return "body_txt"


def classify_non_text(bbox: Tuple[int, int, int, int, int], page_w: int, page_h: int, crop: Image.Image) -> str:
    x, y, w, h, area = bbox
    page_area = page_w * page_h
    aspect = w / max(h, 1)
    box_area = w * h
    if box_area > page_area * 0.25:
        return "figure"
    if h <= max(4, page_h * 0.012) or w <= max(4, page_w * 0.006) or aspect > 12 or aspect < 0.08:
        return "line"
    if box_area > page_area * 0.06:
        return "figure"
    if box_area > page_area * 0.025 and (aspect > 2.0 or aspect < 0.5):
        return "decoration"
    if box_area < page_area * 0.004:
        return "icon"
    # Heuristic: many unique colors suggests photo/figure/chart.
    arr = np.array(crop.convert("RGBA"))
    pix = arr[arr[:, :, 3] > 32, :3]
    if pix.size > 0:
        sample = pix[::max(1, len(pix)//4000)]
        uniq = len(np.unique((sample // 16).astype(np.uint8), axis=0))
        if uniq > 80:
            return "figure"
    return "shape"


def infer_font_size_points(box_h_px: int, source_h: int, slide_h_in: float) -> float:
    # Rough relationship: rendered pixel height to slide points.
    # OCR bbox includes only glyph height, so scale moderately upward.
    points_per_px = (slide_h_in * 72.0) / max(source_h, 1)
    return max(6.0, min(60.0, box_h_px * points_per_px * 1.15))


def normalized_bbox(x: int, y: int, w: int, h: int, source_w: int, source_h: int) -> str:
    return json.dumps([x/source_w, y/source_h, w/source_w, h/source_h])


def emu_coords(x: int, y: int, w: int, h: int, source_w: int, source_h: int, slide_w_emu: int, slide_h_emu: int) -> Tuple[int, int, int, int]:
    return (
        int(x / source_w * slide_w_emu),
        int(y / source_h * slide_h_emu),
        int(w / source_w * slide_w_emu),
        int(h / source_h * slide_h_emu),
    )


def add_background_element(page: PageSource, slide_dir: Path, out_root: Path, z: int, records: List[ElementRecord], slide_w_emu: int, slide_h_emu: int, background_color: Tuple[int, int, int]) -> int:
    bg = Image.new("RGBA", (1, 1), background_color + (255,))
    fname = f"{page.slide_id}_background_01.png"
    fpath = slide_dir / fname
    save_png_no_compress(bg, fpath)
    left, top, width, height = 0, 0, slide_w_emu, slide_h_emu
    rec = ElementRecord(
        slide_id=page.slide_id,
        element_id=f"{page.slide_id}_E{z:04d}",
        file_name=fname,
        relative_path=str(fpath.relative_to(out_root)),
        element_type="background",
        is_text_only="no",
        recognized_text="",
        OCR_confidence=0.0,
        x=0, y=0, width=page.source_width, height=page.source_height,
        z_order=z,
        opacity=1.0,
        rotation=0.0,
        crop_mask_type="full_background",
        transparent_background="no",
        alpha_note="full_background",
        shadow_effect="",
        converted_to_editable_text="n/a",
        conversion_note="solid background representation",
        page_source_type=page.page_source_type,
        source_page_path=page.source_path,
        source_width=page.source_width,
        source_height=page.source_height,
        bbox_pixels=json.dumps([0, 0, page.source_width, page.source_height]),
        bbox_normalized=json.dumps([0, 0, 1, 1]),
        pptx_left_emu=left,
        pptx_top_emu=top,
        pptx_width_emu=width,
        pptx_height_emu=height,
        dominant_color="#%02X%02X%02X" % background_color,
        font_family_guess="",
        font_size_guess=0.0,
        font_color_guess="",
        text_alignment_guess="",
        language_guess="",
        split_confidence=1.0,
        parent_group_id="background",
        needs_manual_review="no",
    )
    records.append(rec)
    return z + 1


def decompose_page(page: PageSource, out_root: Path, args: argparse.Namespace, start_z: int = 1) -> List[ElementRecord]:
    img = Image.open(page.source_path).convert("RGBA")
    rgb_img = img.convert("RGB")
    arr_rgb = np.array(rgb_img)
    page_h, page_w = arr_rgb.shape[:2]
    page.source_width = page_w; page.source_height = page_h
    page_dir = out_root / "split_png_elements" / page.slide_id
    safe_mkdir(page_dir)
    records: List[ElementRecord] = []
    slide_w_emu = int(DEFAULT_SLIDE_W_IN * EMU_PER_INCH)
    slide_h_in = DEFAULT_SLIDE_W_IN * page_h / page_w
    slide_h_emu = int(slide_h_in * EMU_PER_INCH)

    bg_color = estimate_background_color(rgb_img)
    z = start_z
    if not args.no_background:
        z = add_background_element(page, page_dir, out_root, z, records, slide_w_emu, slide_h_emu, bg_color)

    # Text detection first.
    text_boxes = text_boxes_from_ocr(rgb_img, args.ocr_lang, args.ocr, args.ocr_confidence_threshold)
    # Make text box list stable and remove tiny duplicates.
    dedup: List[Dict[str, Any]] = []
    for b in sorted(text_boxes, key=lambda d: (d["y"], d["x"])):
        if any(boxes_iou((b["x"], b["y"], b["w"], b["h"]), (u["x"], u["y"], u["w"], u["h"])) > 0.75 for u in dedup):
            continue
        dedup.append(b)
    text_boxes = dedup

    type_counts: Dict[str, int] = {}
    for tb in text_boxes:
        x, y, w, h = int(tb["x"]), int(tb["y"]), int(tb["w"]), int(tb["h"])
        # Limit to page bounds.
        x = max(0, x); y = max(0, y); w = min(w, page_w - x); h = min(h, page_h - y)
        if w <= 1 or h <= 1:
            continue
        text = tb["text"]
        conf = float(tb["conf"])
        etype = infer_text_type(tb, page_w, page_h, text)
        type_counts[etype] = type_counts.get(etype, 0) + 1
        fname = f"{page.slide_id}_{etype}_{type_counts[etype]:02d}.png"
        crop = crop_text_png(img, (x, y, w, h), bg_color)
        fpath = page_dir / fname
        save_png_no_compress(crop, fpath)
        color = dominant_color_hex(crop)
        crop_mask = infer_crop_mask_type(crop)
        transparent_background, alpha_note = alpha_fields_for(crop, etype, "yes")
        font_size = infer_font_size_points(h, page_h, slide_h_in)
        left, top, width, height = emu_coords(x, y, w, h, page_w, page_h, slide_w_emu, slide_h_emu)
        converted = "yes" if args.editable_text and conf >= args.ocr_confidence_threshold else "no"
        note = "editable text box" if converted == "yes" else ("low_ocr_confidence" if args.editable_text else "editable_text_disabled")
        rec = ElementRecord(
            slide_id=page.slide_id,
            element_id=f"{page.slide_id}_E{z:04d}",
            file_name=fname,
            relative_path=str(fpath.relative_to(out_root)),
            element_type=etype,
            is_text_only="yes",
            recognized_text=text,
            OCR_confidence=round(conf, 2),
            x=x, y=y, width=w, height=h,
            z_order=z,
            opacity=1.0,
            rotation=0.0,
            crop_mask_type=crop_mask,
            transparent_background=transparent_background,
            alpha_note=alpha_note,
            shadow_effect="",
            converted_to_editable_text=converted,
            conversion_note=note,
            page_source_type=page.page_source_type,
            source_page_path=page.source_path,
            source_width=page_w,
            source_height=page_h,
            bbox_pixels=json.dumps([x, y, w, h]),
            bbox_normalized=normalized_bbox(x, y, w, h, page_w, page_h),
            pptx_left_emu=left,
            pptx_top_emu=top,
            pptx_width_emu=width,
            pptx_height_emu=height,
            dominant_color=color,
            font_family_guess=args.default_font_family,
            font_size_guess=round(font_size, 2),
            font_color_guess=color,
            text_alignment_guess="center" if w > page_w * 0.55 and x < page_w * 0.25 else "left",
            language_guess=args.ocr_lang,
            split_confidence=round(min(1.0, conf / 100.0), 3),
            parent_group_id="",
            needs_manual_review="no" if conf >= args.ocr_confidence_threshold else "yes",
        )
        records.append(rec)
        z += 1

    # Non-text segmentation.
    params = granularity_params(args.granularity, page_w * page_h)
    mask = non_background_mask(arr_rgb, bg_color)
    mask_non_text = subtract_text_regions(mask, text_boxes, pad=max(3, page_w // 400))
    kernel = cv2.getStructuringElement(cv2.MORPH_RECT, params["close_kernel"])
    mask_non_text = cv2.morphologyEx(mask_non_text, cv2.MORPH_CLOSE, kernel, iterations=1)
    if params["dilate"]:
        mask_non_text = cv2.dilate(mask_non_text, cv2.getStructuringElement(cv2.MORPH_RECT, (3, 3)), iterations=params["dilate"])
    boxes = connected_boxes(mask_non_text, int(params["min_area"]))
    boxes = merge_boxes(boxes, args.granularity, page_w, page_h)
    # Remove boxes that are too close to text boxes or cover the entire slide.
    filtered: List[Tuple[int, int, int, int, int]] = []
    for b in boxes:
        x, y, w, h, area = b
        if w * h > page_w * page_h * 0.93:
            continue
        if any(boxes_iou((x, y, w, h), (int(t["x"]), int(t["y"]), int(t["w"]), int(t["h"]))) > 0.45 for t in text_boxes):
            continue
        filtered.append(b)

    type_counts_non: Dict[str, int] = {}
    for b in sorted(filtered, key=lambda t: (t[1], t[0], t[2]*t[3])):
        x, y, w, h, area = b
        if w <= 1 or h <= 1:
            continue
        crop = make_alpha_crop(img, mask_non_text, (x, y, w, h))
        if np.array(crop.convert("RGBA"))[:, :, 3].max() <= 10:
            continue
        etype = classify_non_text(b, page_w, page_h, crop)
        type_counts_non[etype] = type_counts_non.get(etype, 0) + 1
        fname = f"{page.slide_id}_{etype}_{type_counts_non[etype]:02d}.png"
        fpath = page_dir / fname
        save_png_no_compress(crop, fpath)
        color = dominant_color_hex(crop)
        crop_mask = infer_crop_mask_type(crop)
        transparent_background, alpha_note = alpha_fields_for(crop, etype, "no")
        left, top, width, height = emu_coords(x, y, w, h, page_w, page_h, slide_w_emu, slide_h_emu)
        needs_review = "yes" if etype in {"mixed", "unknown"} or (w*h > page_w*page_h*0.20) else "no"
        if alpha_note.startswith("transparent_required_but_unreliable"):
            needs_review = "yes"
        rec = ElementRecord(
            slide_id=page.slide_id,
            element_id=f"{page.slide_id}_E{z:04d}",
            file_name=fname,
            relative_path=str(fpath.relative_to(out_root)),
            element_type=etype,
            is_text_only="no",
            recognized_text="",
            OCR_confidence=0.0,
            x=x, y=y, width=w, height=h,
            z_order=z,
            opacity=1.0,
            rotation=0.0,
            crop_mask_type=crop_mask,
            transparent_background=transparent_background,
            alpha_note=alpha_note,
            shadow_effect="unknown" if crop_mask == "irregular_mask" else "",
            converted_to_editable_text="n/a",
            conversion_note="png element",
            page_source_type=page.page_source_type,
            source_page_path=page.source_path,
            source_width=page_w,
            source_height=page_h,
            bbox_pixels=json.dumps([x, y, w, h]),
            bbox_normalized=normalized_bbox(x, y, w, h, page_w, page_h),
            pptx_left_emu=left,
            pptx_top_emu=top,
            pptx_width_emu=width,
            pptx_height_emu=height,
            dominant_color=color,
            font_family_guess="",
            font_size_guess=0.0,
            font_color_guess="",
            text_alignment_guess="",
            language_guess="",
            split_confidence=0.65 if needs_review == "no" else 0.45,
            parent_group_id="",
            needs_manual_review=needs_review,
        )
        records.append(rec)
        z += 1
    # Re-sort by z_order while preserving background first and detected order.
    records.sort(key=lambda r: r.z_order)
    return records


def write_manifest(records: List[ElementRecord], outdir: Path) -> None:
    rows = [asdict(r) for r in records]
    if not rows:
        return
    csv_path = outdir / "visual_elements_manifest.csv"
    with csv_path.open("w", newline="", encoding="utf-8-sig") as f:
        writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
        writer.writeheader(); writer.writerows(rows)
    (outdir / "visual_elements_manifest.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")


def zip_folder(folder: Path, zip_path: Path) -> None:
    if zip_path.exists():
        zip_path.unlink()
    with zipfile.ZipFile(zip_path, "w", compression=zipfile.ZIP_DEFLATED, compresslevel=1) as zf:
        for p in sorted(folder.rglob("*")):
            if p.is_file():
                zf.write(p, p.relative_to(folder.parent))


def create_review_overlays(sources: List[PageSource], records: List[ElementRecord], outdir: Path) -> None:
    review_dir = outdir / "review"
    safe_mkdir(review_dir)
    by_slide: Dict[str, List[ElementRecord]] = {}
    for r in records:
        if r.element_type == "background":
            continue
        by_slide.setdefault(r.slide_id, []).append(r)
    source_map = {s.slide_id: s for s in sources}
    try:
        font = ImageFont.truetype("DejaVuSans.ttf", 16)
    except Exception:
        font = ImageFont.load_default()
    for sid, items in by_slide.items():
        src = source_map[sid]
        img = Image.open(src.source_path).convert("RGB")
        draw = ImageDraw.Draw(img)
        for r in items:
            color = (255, 0, 0) if r.is_text_only == "yes" else (0, 140, 255)
            draw.rectangle([r.x, r.y, r.x + r.width, r.y + r.height], outline=color, width=max(2, img.width // 700))
            label = r.element_id.split("_")[-1] + ":" + r.element_type
            draw.rectangle([r.x, max(0, r.y - 20), r.x + min(180, len(label) * 8), r.y], fill=color)
            draw.text((r.x + 2, max(0, r.y - 18)), label, fill=(255, 255, 255), font=font)
        img.save(review_dir / f"{sid}_detected_elements_overlay.png", quality=95)


def add_textbox(slide: Any, rec: ElementRecord, slide_w: int, slide_h: int, page_w: int, page_h: int) -> None:
    left = int(rec.x / page_w * slide_w)
    top = int(rec.y / page_h * slide_h)
    width = max(1, int(rec.width / page_w * slide_w))
    height = max(1, int(rec.height / page_h * slide_h))
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear(); tf.word_wrap = True; tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    lines = str(rec.recognized_text or "").split("\n")
    for i, line in enumerate(lines if lines else [""]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        align = rec.text_alignment_guess
        p.alignment = PP_ALIGN.CENTER if align == "center" else PP_ALIGN.LEFT
        run_obj = p.add_run()
        run_obj.text = line
        font = run_obj.font
        font.name = rec.font_family_guess or "Arial"
        font.size = Pt(max(4, rec.font_size_guess or 10))
        if rec.font_color_guess and re.match(r"^#[0-9A-Fa-f]{6}$", rec.font_color_guess):
            hexstr = rec.font_color_guess.lstrip("#")
            font.color.rgb = RGBColor(int(hexstr[0:2], 16), int(hexstr[2:4], 16), int(hexstr[4:6], 16))
        else:
            font.color.rgb = RGBColor(0, 0, 0)
        if rec.element_type == "title_txt":
            font.bold = True


def recompose_pptx(sources: List[PageSource], records: List[ElementRecord], outdir: Path, args: argparse.Namespace) -> Path:
    prs = Presentation()
    first = sources[0]
    aspect = first.source_height / first.source_width
    slide_w_in = DEFAULT_SLIDE_W_IN
    slide_h_in = slide_w_in * aspect
    prs.slide_width = int(slide_w_in * EMU_PER_INCH)
    prs.slide_height = int(slide_h_in * EMU_PER_INCH)
    blank = prs.slide_layouts[6]
    source_map = {s.slide_id: s for s in sources}
    by_slide: Dict[str, List[ElementRecord]] = {}
    for r in records:
        by_slide.setdefault(r.slide_id, []).append(r)

    for source in sources:
        slide = prs.slides.add_slide(blank)
        try:
            bg_color = estimate_background_color(Image.open(source.source_path).convert("RGB"))
            slide.background.fill.solid()
            slide.background.fill.fore_color.rgb = RGBColor(*bg_color)
        except Exception:
            pass
        for rec in sorted(by_slide.get(source.slide_id, []), key=lambda r: r.z_order):
            if rec.element_type == "background":
                continue
            if rec.is_text_only == "yes" and rec.converted_to_editable_text == "yes":
                add_textbox(slide, rec, prs.slide_width, prs.slide_height, rec.source_width, rec.source_height)
            else:
                img_path = outdir / rec.relative_path
                if not img_path.exists():
                    continue
                left = int(rec.x / rec.source_width * prs.slide_width)
                top = int(rec.y / rec.source_height * prs.slide_height)
                width = int(rec.width / rec.source_width * prs.slide_width)
                height = int(rec.height / rec.source_height * prs.slide_height)
                slide.shapes.add_picture(str(img_path), left, top, width=max(1, width), height=max(1, height))
    out = outdir / "recomposed_from_elements.pptx"
    prs.save(out)
    return out


def render_pptx_for_quality(pptx_path: Path, outdir: Path, dpi: int = 150) -> List[Path]:
    safe_mkdir(outdir)
    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        try:
            office = libreoffice_command()
        except RuntimeError:
            return []
        proc = run([office, "--headless", "--convert-to", "pdf", "--outdir", str(tmp), str(pptx_path)], check=False)
        if proc.returncode != 0:
            return []
        pdfs = sorted(tmp.glob("*.pdf"))
        if not pdfs:
            return []
        run(["pdftoppm", "-r", str(dpi), "-png", str(pdfs[0]), str(outdir / "slide")], check=False)
    return sorted(outdir.glob("slide-*.png"))


def quality_check(sources: List[PageSource], pptx_path: Path, outdir: Path) -> None:
    qdir = outdir / "quality_report"
    safe_mkdir(qdir)
    render_dir = qdir / "recomposed_render"
    rendered = render_pptx_for_quality(pptx_path, render_dir)
    rows = []
    for idx, source in enumerate(sources):
        if idx >= len(rendered):
            rows.append(dict(slide_id=source.slide_id, pixel_diff_ratio="", mean_absolute_error="", status="render_missing"))
            continue
        orig = Image.open(source.source_path).convert("RGB")
        rec = Image.open(rendered[idx]).convert("RGB")
        orig_resized = orig.resize(rec.size, Image.Resampling.LANCZOS)
        diff = ImageChops.difference(orig_resized, rec)
        arr = np.array(diff)
        gray = arr.mean(axis=2)
        diff_ratio = float(np.mean(gray > 20))
        mae = float(np.mean(gray))
        orig_resized.save(qdir / f"{source.slide_id}_original.png")
        rec.save(qdir / f"{source.slide_id}_recomposed.png")
        # Amplify diff for review.
        amp = np.clip(arr * 4, 0, 255).astype(np.uint8)
        Image.fromarray(amp).save(qdir / f"{source.slide_id}_diff.png")
        rows.append(dict(
            slide_id=source.slide_id,
            pixel_diff_ratio=round(diff_ratio, 5),
            mean_absolute_error=round(mae, 3),
            status="manual_review_required" if diff_ratio > 0.08 else "ok",
        ))
    if rows:
        with (qdir / "quality_report.csv").open("w", newline="", encoding="utf-8-sig") as f:
            writer = csv.DictWriter(f, fieldnames=list(rows[0].keys()))
            writer.writeheader(); writer.writerows(rows)
        (qdir / "quality_report.json").write_text(json.dumps(rows, ensure_ascii=False, indent=2), encoding="utf-8")


def parse_args() -> argparse.Namespace:
    p = argparse.ArgumentParser(description="Decompose image-based PPT/PPTX or screenshots into PNG elements and editable-text recomposed PPTX.")
    p.add_argument("input", type=Path, help="Input PPT/PPTX or image file")
    p.add_argument("--outdir", type=Path, default=Path("output"), help="Output directory")
    p.add_argument("--config", type=Path, help="Optional YAML/JSON config file")
    p.add_argument("--dpi", type=int, default=300, help="DPI for fallback rendering")
    p.add_argument("--granularity", choices=["coarse", "normal", "fine", "ultra"], default="fine")
    p.add_argument("--ocr", action="store_true", help="Enable OCR for text detection")
    p.add_argument("--ocr-lang", default="eng", help="Tesseract language, e.g. chi_sim+eng")
    p.add_argument("--ocr-confidence-threshold", type=float, default=70.0, help="Minimum OCR confidence to convert to editable text")
    p.add_argument("--editable-text", action="store_true", help="Convert reliable text-only PNG elements into editable PPTX text boxes")
    p.add_argument("--default-font-family", default="Arial", help="Font family used for editable text boxes")
    p.add_argument("--review", action="store_true", help="Generate review overlay images with detected element boxes")
    p.add_argument("--quality-check", action="store_true", help="Render recomposed PPTX and generate visual diff report")
    p.add_argument("--no-background", action="store_true", help="Do not export a solid background PNG element")
    p.add_argument("--clean", action="store_true", help="Delete output directory before processing")
    return p.parse_args()


# argparse defaults are needed by config merge.
parser = argparse.ArgumentParser(add_help=False)
parser.add_argument("--dpi", type=int, default=300)
parser.add_argument("--granularity", default="fine")
parser.add_argument("--ocr-lang", default="eng")
parser.add_argument("--ocr-confidence-threshold", type=float, default=70.0)
parser.add_argument("--default-font-family", default="Arial")


def main() -> int:
    args = parse_args()
    cfg = load_config(args.config)
    args = apply_config_defaults(args, cfg)
    input_path = args.input.resolve()
    outdir = args.outdir.resolve()
    if args.clean and outdir.exists():
        shutil.rmtree(outdir)
    safe_mkdir(outdir)

    if not input_path.exists():
        raise FileNotFoundError(input_path)
    if args.ocr and pytesseract is None:
        print("WARNING: pytesseract is not installed; OCR/text conversion will be disabled.", file=sys.stderr)
        args.ocr = False
        args.editable_text = False

    print(f"[1/6] Collecting source pages from {input_path.name}...")
    sources = collect_page_sources(input_path, outdir, args.dpi)
    write_page_source_report(sources, outdir)
    print(f"      pages/images: {len(sources)}")

    print(f"[2/6] Decomposing visual elements, granularity={args.granularity}...")
    records: List[ElementRecord] = []
    for src in sources:
        print(f"      {src.slide_id}: {src.page_source_type} {src.source_width}x{src.source_height}")
        records.extend(decompose_page(src, outdir, args))

    print(f"[3/6] Writing manifest ({len(records)} elements)...")
    write_manifest(records, outdir)

    print("[4/6] Creating element ZIP...")
    zip_folder(outdir / "split_png_elements", outdir / "split_png_elements.zip")

    if args.review:
        print("[5/6] Creating review overlays...")
        create_review_overlays(sources, records, outdir)
    else:
        print("[5/6] Review overlays skipped.")

    print("[6/6] Recomposing PPTX...")
    pptx_path = recompose_pptx(sources, records, outdir, args)

    if args.quality_check:
        print("      Running quality check...")
        quality_check(sources, pptx_path, outdir)

    print("Done.")
    print(f"Output folder: {outdir}")
    print(f"PPTX: {pptx_path}")
    print(f"ZIP: {outdir / 'split_png_elements.zip'}")
    print(f"Manifest CSV: {outdir / 'visual_elements_manifest.csv'}")
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
